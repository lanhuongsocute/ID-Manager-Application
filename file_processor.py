from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
import os

def gan_id_word(file_path, id_text, output_folder):
    doc = Document(file_path)
    for section in doc.sections:
        header = section.header
        if not header.paragraphs:
            header.add_paragraph()
        header.paragraphs[0].text = f"ID: {id_text}"
    new_path = os.path.join(output_folder, f"{id_text}_{os.path.basename(file_path)}")
    doc.save(new_path)
    return new_path

def gan_id_ppt(file_path, id_text, output_folder):
    prs = Presentation(file_path)
    for slide in prs.slides:
        width = prs.slide_width
        height = prs.slide_height
        textbox = slide.shapes.add_textbox(Inches(0.5), height - Inches(0.5), width - Inches(1), Inches(0.3))
        tf = textbox.text_frame
        p = tf.add_paragraph()
        p.text = f"ID: {id_text}"
        p.font.size = Pt(10)
    new_path = os.path.join(output_folder, f"{id_text}_{os.path.basename(file_path)}")
    prs.save(new_path)
    return new_path
