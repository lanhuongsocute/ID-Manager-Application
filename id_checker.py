from docx import Document
from pptx import Presentation

def check_id_word(file_path):
    doc = Document(file_path)
    for section in doc.sections:
        for p in section.header.paragraphs:
            if p.text.strip().startswith("ID:"):
                return p.text.strip()
    for p in doc.paragraphs:
        if p.text.strip().startswith("ID:"):
            return p.text.strip()
    return None

def check_id_ppt(file_path):
    prs = Presentation(file_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                if shape.text.strip().startswith("ID:"):
                    return shape.text.strip()
    return None
